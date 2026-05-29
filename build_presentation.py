#!/usr/bin/env python3
"""
Gera o arquivo PPTX da Jornada do Aluno - BLACKBOX. JIU-JITSU
Para abrir no Google Slides: Upload no Google Drive → Abrir com Google Slides
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ── BRAND COLORS ──────────────────────────────────────────────────────────────
BK  = RGBColor(0x0A, 0x0A, 0x0A)   # blackout (background)
TA  = RGBColor(0x1A, 0x1A, 0x1A)   # tatame (card bg)
GO  = RGBColor(0xE8, 0xB8, 0x4B)   # gold
WH  = RGBColor(0xF5, 0xF5, 0xF0)   # white-belt
MI  = RGBColor(0x66, 0x66, 0x66)   # mid-gray
DM  = RGBColor(0x2A, 0x2A, 0x2A)   # dim

# ── SLIDE SIZE: 16:9 WIDESCREEN ───────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # completely blank

# ── HELPERS ───────────────────────────────────────────────────────────────────
def add_slide():
    sl = prs.slides.add_slide(blank_layout)
    # black background
    bg = sl.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BK
    return sl

def txb(sl, text, x, y, w, h,
        size=14, bold=False, color=WH, align=PP_ALIGN.LEFT,
        font="DM Sans", italic=False):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def txb_lines(sl, lines, x, y, w, h,
              size=12, bold=False, color=WH, align=PP_ALIGN.LEFT,
              font="DM Sans", line_spacing=1.2):
    """Multi-line text box with list of (text, color, bold, size) tuples"""
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, color, bold, size)
        txt, col, bld, sz = item[0], item[1] if len(item)>1 else color, item[2] if len(item)>2 else bold, item[3] if len(item)>3 else size
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(3)
        run = p.add_run()
        run.text = txt
        run.font.name  = font
        run.font.size  = Pt(sz)
        run.font.bold  = bld
        run.font.color.rgb = col
    return tb

def rect(sl, x, y, w, h, fill=TA, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shp = sl.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = line_w
    else:
        shp.line.fill.background()  # no line
    return shp

def gold_bar(sl, x, y, w=Inches(0.4), h=Inches(0.025)):
    """Thin gold divider line"""
    rect(sl, x, y, w, h, fill=GO)

def progress_bar(sl, slide_num, total=11):
    """Top progress bar"""
    pct = (slide_num - 1) / (total - 1) if total > 1 else 0
    rect(sl, Inches(0), Inches(0), Inches(13.33 * pct), Inches(0.04), fill=GO)

def tag_label(sl, text, x, y):
    """Small gold uppercase tracking label"""
    txb(sl, text, x, y, Inches(8), Inches(0.3),
        size=9, color=GO, font="Courier New")

def big_title(sl, text, x, y, size=52):
    """Bebas-style headline (using Arial Black as fallback)"""
    txb(sl, text, x, y, Inches(9), Inches(1.4),
        size=size, bold=True, color=WH, font="Arial Black")

def card(sl, x, y, w, h, gold_left=False, gold_bg=False):
    """Card with dark tatame background"""
    bg_color = RGBColor(0x14, 0x12, 0x08) if gold_bg else TA
    shp = rect(sl, x, y, w, h, fill=bg_color,
               line=GO if gold_bg else RGBColor(0x28, 0x28, 0x28),
               line_w=Pt(1))
    if gold_left:
        rect(sl, x, y, Inches(0.04), h, fill=GO)
    return shp

def step_circle(sl, x, y, text, done=True, hot=False):
    """Flow step circle"""
    r = Inches(0.45)
    c = GO if hot else (RGBColor(0x44, 0x40, 0x28) if done else DM)
    shp = sl.shapes.add_shape(9, x, y, r, r)  # oval
    shp.fill.solid()
    shp.fill.fore_color.rgb = DM
    shp.line.color.rgb = c
    shp.line.width = Pt(1.5)
    tf = shp.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Arial Black"
    run.font.size = Pt(9 if len(text) > 2 else 11)
    run.font.bold = True
    run.font.color.rgb = c
    return shp

def arrow_right(sl, x, y, w=Inches(0.22)):
    """Small right arrow"""
    rect(sl, x, y + Inches(0.2), w, Inches(0.02), fill=MI)

def belt_bar_h(sl, x, y, w, h, color, border_gold=False):
    bdr = GO if border_gold else None
    rect(sl, x, y, w, h, fill=color, line=bdr, line_w=Pt(0.75))

def highlight_box(sl, x, y, w, h, label, value, desc):
    card(sl, x, y, w, h, gold_bg=True)
    txb(sl, label, x + Inches(0.18), y + Inches(0.14), w, Inches(0.22),
        size=7, color=RGBColor(0xCC, 0xA0, 0x3A), font="Courier New")
    txb(sl, value, x + Inches(0.18), y + Inches(0.32), w, Inches(0.5),
        size=28, bold=True, color=GO, font="Arial Black")
    txb(sl, desc, x + Inches(0.18), y + Inches(0.82), w - Inches(0.2), Inches(0.5),
        size=9, color=MI, font="DM Sans")

# ── BELT COLORS ────────────────────────────────────────────────────────────────
BELTS = [
    (RGBColor(0xF5,0xF5,0xF0), "FAIXA BRANCA",  "Fundamentos · 2 anos mín.",    False),
    (RGBColor(0x25,0x63,0xEB), "FAIXA AZUL",    "Consolidação técnica · 2 anos mín.", False),
    (RGBColor(0x7C,0x3A,0xED), "FAIXA ROXA",    "Jogo próprio · 1,5 ano mín.",  False),
    (RGBColor(0x92,0x40,0x0E), "FAIXA MARROM",  "Refinamento · 1 ano mín.",     False),
    (RGBColor(0x0A,0x0A,0x0A), "FAIXA PRETA",   "Maestria · Legado no esporte", True),
]

FLOW_STEPS = [
    ("01","WA",       "WhatsApp\nContato"),
    ("02","Site",     "Site\nBlackbox"),
    ("03","Form",     "Formulário\nAgend."),
    ("04","Aula",     "Aula\nExp."),
    ("05","Base",     "Cadastro\nBaseBJJ"),
    ("06","Pag",      "Pagamento\nMatríc."),
    ("07","BJJ",      "Acesso\nAcademia"),
    ("08","Grad",     "Graduação\nEvolução"),
]

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 0 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 1)

# Belt bars decoration
bx = Inches(5.0)
for i in range(4):
    rect(sl, bx + Inches(i * 0.09), Inches(1.2), Inches(0.04), Inches(0.35), fill=RGBColor(0x50,0x45,0x1A))
rect(sl, bx + Inches(0.44), Inches(1.22), Inches(0.22), Inches(0.22),
     fill=BK, line=RGBColor(0x50,0x45,0x1A), line_w=Pt(1))

# Main wordmark
txb(sl, "BLACKBOX.", Inches(2.2), Inches(1.8), Inches(9), Inches(2.2),
    size=96, bold=True, color=WH, font="Arial Black", align=PP_ALIGN.CENTER)
# gold dot on the period — overlay
txb(sl, ".", Inches(9.25), Inches(1.8), Inches(1.5), Inches(2.2),
    size=96, bold=True, color=GO, font="Arial Black")

txb(sl, "JORNADA DO ALUNO", Inches(2.0), Inches(4.0), Inches(9.33), Inches(0.9),
    size=36, bold=True, color=RGBColor(0xAA,0xAA,0xA0), font="Arial Black", align=PP_ALIGN.CENTER)

# Divider
rect(sl, Inches(4.2), Inches(4.95), Inches(5.0), Inches(0.02), fill=DM)

txb(sl, "Do primeiro contato à faixa preta", Inches(2), Inches(5.1), Inches(9.33), Inches(0.4),
    size=11, color=GO, font="Courier New", align=PP_ALIGN.CENTER)
txb(sl, "APRESENTAÇÃO INTERNA · 2025", Inches(2), Inches(5.55), Inches(9.33), Inches(0.4),
    size=9, color=MI, font="Courier New", align=PP_ALIGN.CENTER)

# Pills
pill_texts = ["3 UNIDADES", "5 MODALIDADES", "APP BASEBJJ", "SITE BLACKBOX"]
pill_x = [Inches(1.8), Inches(3.9), Inches(6.0), Inches(7.9)]
for i, (pt, px) in enumerate(zip(pill_texts, pill_x)):
    rect(sl, px, Inches(6.15), Inches(1.8), Inches(0.28), fill=RGBColor(0x18,0x14,0x05), line=GO, line_w=Pt(0.5))
    txb(sl, pt, px, Inches(6.15), Inches(1.8), Inches(0.28),
        size=7, color=GO, font="Courier New", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — VISÃO GERAL
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 2)
tag_label(sl, "VISÃO GERAL · 8 ETAPAS", Inches(0.55), Inches(0.55))
big_title(sl, "JORNADA COMPLETA", Inches(0.55), Inches(0.85), size=44)
gold_bar(sl, Inches(0.55), Inches(1.65))

txb(sl, "Do primeiro contato no WhatsApp ao tatame — com acompanhamento digital via BaseBJJ e gestão ativa de graduações pelo professor.",
    Inches(0.55), Inches(1.75), Inches(9), Inches(0.5), size=11, color=MI)

# Flow chart
step_w = Inches(1.4)
gap    = Inches(0.08)
start_x = Inches(0.35)
flow_y  = Inches(2.45)
for i, (sid, short, lbl) in enumerate(FLOW_STEPS):
    sx = start_x + i * (step_w + gap)
    step_circle(sl, sx + Inches(0.47), flow_y, short, done=True)
    txb(sl, lbl, sx, flow_y + Inches(0.55), step_w, Inches(0.5),
        size=8, color=RGBColor(0x88,0x84,0x60), font="Courier New", align=PP_ALIGN.CENTER)
    if i < len(FLOW_STEPS) - 1:
        arrow_right(sl, sx + step_w - Inches(0.05), flow_y + Inches(0.18))

rect(sl, Inches(0.35), Inches(3.55), Inches(12.7), Inches(0.02), fill=DM)

# KPI boxes
kpis = [
    ("CONVERSÃO ESPERADA",  "80%",    "Alunos que fazem a aula experimental\ne se matriculam"),
    ("TEMPO DE RESPOSTA",   "< 5min", "No WhatsApp aumenta 8× a\nchance de converter o lead"),
    ("GESTÃO DIGITAL",      "BaseBJJ","Cadastro, pagamentos e\ngraduações integrados"),
    ("PRÉ-VENDA DIGITAL",   "Site",   "Aluno chega motivado e\ninformado antes da 1ª aula"),
]
kpi_w = Inches(3.0)
kpi_h = Inches(2.5)
for i, (lbl, val, desc) in enumerate(kpis):
    kx = Inches(0.35) + i * (kpi_w + Inches(0.14))
    highlight_box(sl, kx, Inches(3.7), kpi_w, kpi_h, lbl, val, desc)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PRIMEIRO CONTATO
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 3)
# Watermark number
txb(sl, "01", Inches(9.5), Inches(0.3), Inches(3.5), Inches(2.5),
    size=150, bold=True, color=RGBColor(0x15,0x12,0x05), font="Arial Black")

tag_label(sl, "ETAPA 01 · CAPTAÇÃO", Inches(0.55), Inches(0.55))
big_title(sl, "PRIMEIRO CONTATO", Inches(0.55), Inches(0.85), size=46)
gold_bar(sl, Inches(0.55), Inches(1.65))

# Left col — bullets
txb(sl, "CANAL DE ENTRADA", Inches(0.55), Inches(1.82), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")

bullets_left = [
    "WhatsApp — canal principal. Resposta em menos de 5 min = 8× mais conversão.",
    "Instagram — DMs e stories com link para o site. Conteúdo de treinos gera interesse.",
    "Indicação de alunos — maior taxa de conversão e menor custo de aquisição.",
    "Presença física — passante ou familiar de aluno na recepção da unidade.",
    "Tráfego pago — Meta Ads e Google Ads. Segmentação local em Cajamar e Santana de Parnaíba.",
]
by = Inches(2.1)
for b in bullets_left:
    rect(sl, Inches(0.55), by + Inches(0.07), Inches(0.07), Inches(0.07), fill=GO)
    txb(sl, b, Inches(0.75), by, Inches(5.2), Inches(0.38), size=10, color=MI)
    by += Inches(0.42)

rect(sl, Inches(0.55), by, Inches(5.5), Inches(0.015), fill=DM)
by += Inches(0.12)

# Template card
card(sl, Inches(0.55), by, Inches(5.5), Inches(1.5), gold_left=True, gold_bg=True)
txb(sl, "TEMPLATE DE BOAS-VINDAS", Inches(0.75), by + Inches(0.12), Inches(5.2), Inches(0.25),
    size=7, color=GO, font="Courier New")
template = (
    '"Olá [nome]! Seja bem-vindo à Blackbox Jiu-Jitsu 🥋\n'
    'Aqui está nosso site para conhecer a academia:\n'
    'blackbox.com.br\n'
    'Quando gostaria de fazer sua aula experimental gratuita?"'
)
txb(sl, template, Inches(0.75), by + Inches(0.38), Inches(5.2), Inches(1.0),
    size=10, color=WH, italic=True)

# Right col — sources
txb(sl, "FONTES DE LEADS", Inches(6.5), Inches(1.82), Inches(6.3), Inches(0.28),
    size=8, color=MI, font="Courier New")

sources = [
    ("A", "Instagram / Redes Sociais",
     "Stories com link para o site. Conteúdo de treinos gera interesse orgânico."),
    ("B", "Indicação de Alunos",
     "Programa: aluno que traz amigo ganha benefício. Maior LTV, menor CAC."),
    ("C", "Tráfego Pago",
     "Meta Ads e Google Ads com segmentação local. CTA direto para WhatsApp."),
    ("D", "SEO / Google Meu Negócio",
     'Busca "academia jiu-jitsu cajamar". Perfis verificados com avaliações.'),
]
sy = Inches(2.1)
for (letter, title, desc) in sources:
    txb(sl, letter, Inches(6.5), sy, Inches(0.3), Inches(0.35),
        size=20, bold=True, color=RGBColor(0x44,0x3A,0x18), font="Arial Black")
    rect(sl, Inches(6.9), sy + Inches(0.04), Inches(5.8), Inches(0.015), fill=DM)
    txb(sl, title, Inches(6.9), sy + Inches(0.06), Inches(5.8), Inches(0.28),
        size=11, bold=True, color=WH, font="DM Sans")
    txb(sl, desc, Inches(6.9), sy + Inches(0.32), Inches(5.8), Inches(0.32),
        size=10, color=MI)
    sy += Inches(0.88)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ENVIO DO SITE
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 4)
txb(sl, "02", Inches(9.5), Inches(0.3), Inches(3.5), Inches(2.5),
    size=150, bold=True, color=RGBColor(0x15,0x12,0x05), font="Arial Black")

tag_label(sl, "ETAPA 02 · PRÉ-VENDA DIGITAL", Inches(0.55), Inches(0.55))
big_title(sl, "ENVIO DO SITE", Inches(0.55), Inches(0.85), size=46)
gold_bar(sl, Inches(0.55), Inches(1.65))

txb(sl, "O SITE COMO FERRAMENTA DE VENDAS", Inches(0.55), Inches(1.82), Inches(5.8), Inches(0.28),
    size=8, color=MI, font="Courier New")

site_bullets = [
    "Link enviado imediatamente no primeiro contato no WhatsApp.",
    "O aluno visita o site antes da aula experimental — chega motivado.",
    "Site mobile-first: projetado para ser lido no celular.",
    "Cada página reforça a identidade da marca e a autoridade Blackbox.",
    "CTAs em todas as páginas levam ao formulário de agendamento.",
]
by = Inches(2.1)
for b in site_bullets:
    rect(sl, Inches(0.55), by + Inches(0.07), Inches(0.07), Inches(0.07), fill=GO)
    txb(sl, b, Inches(0.75), by, Inches(5.2), Inches(0.38), size=10, color=MI)
    by += Inches(0.42)

# Page pills
txb(sl, "PÁGINAS DO SITE", Inches(0.55), by + Inches(0.1), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")
pages = [("INÍCIO", True), ("A BLACKBOX", False), ("MODALIDADES", False),
         ("METODOLOGIA", False), ("PLANOS", True), ("HORÁRIOS", False),
         ("UNIDADES", False), ("CONTATO ←", True)]
pill_colors = [(GO, RGBColor(0x18,0x14,0x05)) if p[1] else (MI, DM) for p in pages]
cols, row_y = 4, by + Inches(0.42)
for i, ((pname, _), (pc, pbg)) in enumerate(zip(pages, pill_colors)):
    px = Inches(0.55) + (i % cols) * Inches(1.35)
    py = row_y + (i // cols) * Inches(0.35)
    rect(sl, px, py, Inches(1.2), Inches(0.25), fill=pbg,
         line=pc, line_w=Pt(0.5))
    txb(sl, pname, px, py, Inches(1.2), Inches(0.25), size=7, color=pc,
        font="Courier New", align=PP_ALIGN.CENTER)

# Right col — navigation flow
txb(sl, "JORNADA DE NAVEGAÇÃO", Inches(6.5), Inches(1.82), Inches(6.3), Inches(0.28),
    size=8, color=MI, font="Courier New")

nav_steps = [
    ("1", "Início — Impacto Visual",
     "Hero BLACKBOX., 3 unidades em destaque e CTA Aula Experimental Grátis."),
    ("2", "Modalidades — O que me interessa?",
     "5 programas com descrição e perfil: BJJ Gi, No-Gi, Feminino, Kids, Defesa."),
    ("3", "Planos — Quanto custa?",
     "Kids, Adulto e Família. FAQ com dúvidas comuns. CTA para WhatsApp por plano."),
    ("4", "Unidades — Qual fica perto?",
     "3 unidades com modalidades, endereços e CTA específico por unidade."),
    ("5", "Contato — Agendar aula",
     "Formulário com unidade + modalidade → WhatsApp com mensagem pré-formatada."),
]
sy = Inches(2.1)
for (num, title, desc) in nav_steps:
    txb(sl, num, Inches(6.5), sy, Inches(0.3), Inches(0.35),
        size=20, bold=True, color=RGBColor(0x44,0x3A,0x18), font="Arial Black")
    rect(sl, Inches(6.9), sy + Inches(0.04), Inches(6.0), Inches(0.015), fill=DM)
    txb(sl, title, Inches(6.9), sy + Inches(0.06), Inches(6.0), Inches(0.28),
        size=11, bold=True, color=WH)
    txb(sl, desc, Inches(6.9), sy + Inches(0.32), Inches(6.0), Inches(0.32),
        size=10, color=MI)
    sy += Inches(0.88)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — FORMULÁRIO & AGENDAMENTO
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 5)
txb(sl, "03", Inches(9.5), Inches(0.3), Inches(3.5), Inches(2.5),
    size=150, bold=True, color=RGBColor(0x15,0x12,0x05), font="Arial Black")

tag_label(sl, "ETAPA 03 · CONVERSÃO", Inches(0.55), Inches(0.55))
big_title(sl, "FORMULÁRIO & AGENDAMENTO", Inches(0.55), Inches(0.85), size=38)
gold_bar(sl, Inches(0.55), Inches(1.65))

# Left
txb(sl, "COMO FUNCIONA", Inches(0.55), Inches(1.82), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")
form_bullets = [
    "Prospect preenche: nome, telefone, unidade e modalidade.",
    "Botão abre o WhatsApp com mensagem pré-formatada — sem atrito.",
    "Equipe recebe mensagem estruturada e confirma o horário.",
    "Aula experimental: gratuita, sem kimono, sem compromisso.",
    "Confirmação via WhatsApp com horário, endereço e orientações.",
]
by = Inches(2.1)
for b in form_bullets:
    rect(sl, Inches(0.55), by + Inches(0.07), Inches(0.07), Inches(0.07), fill=GO)
    txb(sl, b, Inches(0.75), by, Inches(5.2), Inches(0.38), size=10, color=MI)
    by += Inches(0.42)

rect(sl, Inches(0.55), by + Inches(0.1), Inches(5.5), Inches(0.015), fill=DM)

# Template card
card(sl, Inches(0.55), by + Inches(0.22), Inches(5.5), Inches(1.6), gold_left=True, gold_bg=True)
txb(sl, "MENSAGEM ENVIADA AO PROFESSOR", Inches(0.75), by + Inches(0.34), Inches(5.2), Inches(0.25),
    size=7, color=GO, font="Courier New")
msg = ('"Olá! Me chamo João Silva.\n'
       'Quero agendar uma aula experimental na Blackbox.\n'
       'Unidade: Portais — Cajamar (Matriz)\n'
       'Modalidade: BJJ Adultos Gi\n'
       'Meu contato: (11) 99999-9999"')
txb(sl, msg, Inches(0.75), by + Inches(0.6), Inches(5.2), Inches(1.1),
    size=10, color=WH, italic=True)

# Right — Form mockup
card(sl, Inches(6.8), Inches(1.75), Inches(6.0), Inches(5.4))
rect(sl, Inches(6.8), Inches(1.75), Inches(6.0), Inches(0.04), fill=GO)  # top bar
txb(sl, "blackbox.com.br/contato", Inches(7.0), Inches(1.88), Inches(5.5), Inches(0.28),
    size=8, color=GO, font="Courier New")

form_fields = [
    ("NOME", "João Silva"),
    ("TELEFONE / WHATSAPP", "(11) 99999-9999"),
    ("UNIDADE", "Portais — Cajamar (Matriz)"),
    ("MODALIDADE DE INTERESSE", "BJJ Adultos Gi"),
]
fy = Inches(2.22)
for (lbl, val) in form_fields:
    txb(sl, lbl, Inches(7.0), fy, Inches(5.5), Inches(0.22), size=7, color=MI, font="Courier New")
    rect(sl, Inches(7.0), fy + Inches(0.22), Inches(5.6), Inches(0.32), fill=RGBColor(0x20,0x20,0x20))
    txb(sl, val, Inches(7.15), fy + Inches(0.24), Inches(5.4), Inches(0.28), size=10, color=WH)
    fy += Inches(0.7)

rect(sl, Inches(7.0), fy + Inches(0.1), Inches(5.6), Inches(0.48), fill=GO)
txb(sl, "QUERO COMEÇAR →", Inches(7.0), fy + Inches(0.1), Inches(5.6), Inches(0.48),
    size=18, bold=True, color=BK, font="Arial Black", align=PP_ALIGN.CENTER)
txb(sl, "Você será redirecionado para o WhatsApp",
    Inches(7.0), fy + Inches(0.62), Inches(5.6), Inches(0.28),
    size=8, color=MI, font="Courier New", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AULA EXPERIMENTAL
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 6)
txb(sl, "04", Inches(9.5), Inches(0.3), Inches(3.5), Inches(2.5),
    size=150, bold=True, color=RGBColor(0x15,0x12,0x05), font="Arial Black")

tag_label(sl, "ETAPA 04 · O MOMENTO DECISIVO", Inches(0.55), Inches(0.55))
big_title(sl, "AULA EXPERIMENTAL", Inches(0.55), Inches(0.85), size=46)
gold_bar(sl, Inches(0.55), Inches(1.65))

txb(sl, "PROTOCOLO DE RECEPÇÃO", Inches(0.55), Inches(1.82), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")

protocol = [
    ("1", "Acolhimento imediato",
     "Recepcionar pelo nome. Apresentar as instalações. Eliminar o nervosismo."),
    ("2", "Apresentação à equipe",
     "Apresentar o professor e a turma. Criar sentimento de pertencimento."),
    ("3", "Participação na aula",
     "Aquecimento → técnica → drilling → sparring controlado para iniciantes."),
    ("4", "Conversa pós-treino",
     "Escuta ativa: 'O que você achou?' Apresentar o plano ideal ao perfil."),
    ("5", "Fechar na hora",
     "Ter o app BaseBJJ pronto. A decisão é emocional — facilitar o processo."),
]
sy = Inches(2.1)
for (num, title, desc) in protocol:
    txb(sl, num, Inches(0.55), sy, Inches(0.3), Inches(0.35),
        size=20, bold=True, color=RGBColor(0x44,0x3A,0x18), font="Arial Black")
    rect(sl, Inches(0.95), sy + Inches(0.06), Inches(5.0), Inches(0.015), fill=DM)
    txb(sl, title, Inches(0.95), sy + Inches(0.08), Inches(5.0), Inches(0.28),
        size=11, bold=True, color=WH)
    txb(sl, desc, Inches(0.95), sy + Inches(0.34), Inches(5.0), Inches(0.3),
        size=10, color=MI)
    sy += Inches(0.84)

# Right
highlight_box(sl, Inches(6.8), Inches(1.75), Inches(5.9), Inches(1.3),
              "TAXA DE CONVERSÃO", "80%",
              "Das pessoas que fazem a aula\nexperimental se matriculam")

rect(sl, Inches(6.8), Inches(3.1), Inches(5.9), Inches(0.015), fill=DM)

txb(sl, "OBJEÇÕES MAIS COMUNS", Inches(6.8), Inches(3.18), Inches(5.9), Inches(0.28),
    size=8, color=MI, font="Courier New")
objections = [
    '"Vou me sentir fora de lugar?" → Ambiente acolhedor. Todos foram faixa branca.',
    '"Quanto custa?" → Planos claros, sem surpresas. Mostrar custo-benefício.',
    '"Tenho tempo?" → Grade ampla, 3 unidades, reposição de aulas.',
    '"Preciso comprar kimono?" → Não. Pode começar com roupa de treino.',
]
oy = Inches(3.46)
for obj in objections:
    rect(sl, Inches(6.8), oy + Inches(0.07), Inches(0.07), Inches(0.07), fill=GO)
    txb(sl, obj, Inches(7.0), oy, Inches(5.7), Inches(0.38), size=10, color=MI)
    oy += Inches(0.42)

card(sl, Inches(6.8), oy + Inches(0.12), Inches(5.9), Inches(0.9), gold_left=True, gold_bg=True)
txb(sl, "FRASE DE FECHAMENTO", Inches(7.0), oy + Inches(0.24), Inches(5.7), Inches(0.25),
    size=7, color=GO, font="Courier New")
txb(sl, '"Você gostou? Que tal começar semana que vem? Posso mostrar o cadastro agora — leva 3 minutos."',
    Inches(7.0), oy + Inches(0.48), Inches(5.7), Inches(0.45),
    size=10, color=WH, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — CADASTRO BASEBJJ
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 7)
txb(sl, "05", Inches(9.5), Inches(0.3), Inches(3.5), Inches(2.5),
    size=150, bold=True, color=RGBColor(0x15,0x12,0x05), font="Arial Black")

tag_label(sl, "ETAPA 05 · CADASTRO DIGITAL", Inches(0.55), Inches(0.55))
big_title(sl, "CADASTRO — APP BASEBJJ", Inches(0.55), Inches(0.85), size=40)
gold_bar(sl, Inches(0.55), Inches(1.65))

txb(sl, "O QUE É O BASEBJJ", Inches(0.55), Inches(1.82), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")
basebjj_info = [
    "basebjj.com.br — plataforma de gestão para academias de Jiu-Jitsu.",
    "Gerencia alunos, matrículas, pagamentos, frequência e graduações.",
    "Professor tem acesso via painel web. Aluno acessa histórico pelo app.",
    "Notificações automáticas de vencimento reduzem inadimplência.",
]
by = Inches(2.1)
for b in basebjj_info:
    rect(sl, Inches(0.55), by + Inches(0.07), Inches(0.07), Inches(0.07), fill=GO)
    txb(sl, b, Inches(0.75), by, Inches(5.2), Inches(0.38), size=10, color=MI)
    by += Inches(0.42)

rect(sl, Inches(0.55), by + Inches(0.08), Inches(5.5), Inches(0.015), fill=DM)
txb(sl, "DADOS COLETADOS NO CADASTRO", Inches(0.55), by + Inches(0.2), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")
data_fields = [
    "Pessoais: nome completo, CPF, data de nascimento, endereço.",
    "Contato: e-mail, celular, contato de emergência.",
    "Graduação: faixa atual e número de graus.",
    "Plano escolhido e unidade principal de treino.",
    "Termo de responsabilidade digital assinado eletronicamente.",
    "Foto do aluno para identificação no sistema.",
]
dy = by + Inches(0.5)
for d in data_fields:
    rect(sl, Inches(0.55), dy + Inches(0.07), Inches(0.07), Inches(0.07), fill=GO)
    txb(sl, d, Inches(0.75), dy, Inches(5.2), Inches(0.36), size=10, color=MI)
    dy += Inches(0.4)

# Right — BaseBJJ mockup
card(sl, Inches(6.8), Inches(1.75), Inches(6.0), Inches(5.4))
rect(sl, Inches(6.8), Inches(1.75), Inches(6.0), Inches(0.04), fill=GO)
txb(sl, "App BaseBJJ · Perfil do Aluno", Inches(7.0), Inches(1.88), Inches(5.5), Inches(0.28),
    size=8, color=GO, font="Courier New")

# Avatar circle
avatar = sl.shapes.add_shape(9, Inches(7.0), Inches(2.22), Inches(0.55), Inches(0.55))
avatar.fill.solid(); avatar.fill.fore_color.rgb = DM
avatar.line.color.rgb = RGBColor(0x44,0x3A,0x18); avatar.line.width = Pt(1)
txb(sl, "JS", Inches(7.0), Inches(2.26), Inches(0.55), Inches(0.44),
    size=13, bold=True, color=WH, font="Arial Black", align=PP_ALIGN.CENTER)

txb(sl, "JOÃO SILVA", Inches(7.65), Inches(2.24), Inches(5.0), Inches(0.3),
    size=14, bold=True, color=WH, font="Arial Black")
rect(sl, Inches(7.65), Inches(2.58), Inches(1.5), Inches(0.24),
     fill=RGBColor(0x18,0x14,0x05), line=GO, line_w=Pt(0.5))
txb(sl, "FAIXA BRANCA · 0 GRAUS", Inches(7.65), Inches(2.58), Inches(1.5), Inches(0.24),
    size=7, color=GO, font="Courier New", align=PP_ALIGN.CENTER)

rect(sl, Inches(6.9), Inches(2.96), Inches(5.8), Inches(0.015), fill=DM)

rows = [
    ("PLANO",      "Adulto — Ilimitado"),
    ("UNIDADE",    "Portais · Matriz"),
    ("STATUS",     "● Ativo"),
    ("VENCIMENTO", "Dia 05 / mês"),
    ("FREQUÊNCIA", "0 aulas"),
    ("MODALIDADE", "BJJ Adultos Gi"),
]
ry = Inches(3.0)
for lbl, val in rows:
    txb(sl, lbl, Inches(6.95), ry, Inches(2.2), Inches(0.28), size=8, color=MI, font="Courier New")
    col = GO if val.startswith("●") else WH
    txb(sl, val, Inches(9.2), ry, Inches(3.5), Inches(0.28), size=9, color=col, font="Courier New")
    rect(sl, Inches(6.9), ry + Inches(0.3), Inches(5.8), Inches(0.01), fill=DM)
    ry += Inches(0.33)

rect(sl, Inches(6.9), ry + Inches(0.05), Inches(5.8), Inches(0.32),
     fill=RGBColor(0x14,0x12,0x05), line=GO, line_w=Pt(0.5))
txb(sl, "Termo de Responsabilidade · Assinado ✓",
    Inches(6.9), ry + Inches(0.05), Inches(5.8), Inches(0.32),
    size=8, color=GO, font="Courier New", align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PAGAMENTO
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 8)
txb(sl, "06", Inches(9.5), Inches(0.3), Inches(3.5), Inches(2.5),
    size=150, bold=True, color=RGBColor(0x15,0x12,0x05), font="Arial Black")

tag_label(sl, "ETAPA 06 · ATIVAÇÃO FINANCEIRA", Inches(0.55), Inches(0.55))
big_title(sl, "PAGAMENTO & MATRÍCULA", Inches(0.55), Inches(0.85), size=42)
gold_bar(sl, Inches(0.55), Inches(1.65))

txb(sl, "PROCESSO VIA BASEBJJ", Inches(0.55), Inches(1.82), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")
payment_steps = [
    ("1", "Taxa de Matrícula",
     "Cobrada uma única vez. Pode ser zerada em campanhas — a critério do mestre."),
    ("2", "Primeira Mensalidade",
     "Paga no ato do cadastro. Pro-rata ou mês cheio. PIX, crédito ou débito."),
    ("3", "Notificações Automáticas",
     "BaseBJJ envia lembretes de vencimento: D-3, D-0 e D+3 via WhatsApp."),
    ("4", "Gestão de Inadimplência",
     "Aluno em atraso sinalizado. Professor notificado. Protocolo de recuperação acionado."),
    ("5", "Aluno Ativado",
     "Status Ativo no BaseBJJ. Professor notificado: Novo aluno matriculado."),
]
sy = Inches(2.1)
for num, title, desc in payment_steps:
    txb(sl, num, Inches(0.55), sy, Inches(0.3), Inches(0.35),
        size=20, bold=True, color=RGBColor(0x44,0x3A,0x18), font="Arial Black")
    rect(sl, Inches(0.95), sy + Inches(0.06), Inches(5.0), Inches(0.015), fill=DM)
    txb(sl, title, Inches(0.95), sy + Inches(0.08), Inches(5.0), Inches(0.28),
        size=11, bold=True, color=WH)
    txb(sl, desc, Inches(0.95), sy + Inches(0.34), Inches(5.0), Inches(0.3),
        size=10, color=MI)
    sy += Inches(0.8)

# Right — KPI grid
kpi2 = [
    ("FORMAS ACEITAS",   "PIX",    "Crédito · Débito · Boleto"),
    ("RECORRÊNCIA",      "Mensal", "Gestão completa pelo BaseBJJ"),
    ("RÉGUA AUTOMÁTICA", "Auto",   "D-3 · D-0 · D+3 via WhatsApp"),
    ("RECIBO DIGITAL",   "Email",  "Enviado após cada pagamento"),
]
gx = [Inches(6.8), Inches(9.75)]
gy = [Inches(1.75), Inches(3.7)]
for i, (lbl, val, desc) in enumerate(kpi2):
    hx = gx[i % 2]
    hy = gy[i // 2]
    highlight_box(sl, hx, hy, Inches(2.85), Inches(1.75), lbl, val, desc)

card(sl, Inches(6.8), Inches(5.55), Inches(6.0), Inches(1.1), gold_left=True, gold_bg=True)
txb(sl, "APÓS O PAGAMENTO", Inches(7.0), Inches(5.67), Inches(5.7), Inches(0.25),
    size=7, color=GO, font="Courier New")
after_payment = [
    "Aluno marcado como Ativo no BaseBJJ.",
    "Acesso liberado ao registro de frequência e histórico.",
    "Professor recebe notificação de novo aluno matriculado.",
]
ay = Inches(5.95)
for a in after_payment:
    rect(sl, Inches(7.0), ay + Inches(0.06), Inches(0.06), Inches(0.06), fill=GO)
    txb(sl, a, Inches(7.16), ay, Inches(5.5), Inches(0.28), size=10, color=WH)
    ay += Inches(0.3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ACESSO À ACADEMIA
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 9)
txb(sl, "07", Inches(9.5), Inches(0.3), Inches(3.5), Inches(2.5),
    size=150, bold=True, color=RGBColor(0x15,0x12,0x05), font="Arial Black")

tag_label(sl, "ETAPA 07 · INTEGRAÇÃO E ROTINA", Inches(0.55), Inches(0.55))
big_title(sl, "ACESSO À ACADEMIA", Inches(0.55), Inches(0.85), size=44)
gold_bar(sl, Inches(0.55), Inches(1.65))

txb(sl, "ONBOARDING — PRIMEIROS 30 DIAS", Inches(0.55), Inches(1.82), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")
onboarding = [
    "Semana 1: Apresentação formal à turma. Dinâmica das aulas explicada.",
    "Semana 2: Primeiro kimono. Posições básicas, defesa, movimentação.",
    "Semana 3: Sparring supervisionado com parceiros de nível similar.",
    "Semana 4: Check-in com o professor. Aluno percebe primeiros progressos.",
]
by = Inches(2.1)
for b in onboarding:
    rect(sl, Inches(0.55), by + Inches(0.07), Inches(0.07), Inches(0.07), fill=GO)
    txb(sl, b, Inches(0.75), by, Inches(5.2), Inches(0.38), size=10, color=MI)
    by += Inches(0.42)

rect(sl, Inches(0.55), by + Inches(0.08), Inches(5.5), Inches(0.015), fill=DM)
txb(sl, "FERRAMENTAS DE ENGAJAMENTO", Inches(0.55), by + Inches(0.2), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")
engagement = [
    "Grupo de WhatsApp da turma — eventos, horários, motivação.",
    "Frequência no BaseBJJ — aluno acompanha evolução em tempo real.",
    "Instagram da academia — aluno marcado em posts cria vínculo social.",
    "3 unidades acessíveis — flexibilidade aumenta a retenção.",
]
ey = by + Inches(0.5)
for e in engagement:
    rect(sl, Inches(0.55), ey + Inches(0.07), Inches(0.07), Inches(0.07), fill=GO)
    txb(sl, e, Inches(0.75), ey, Inches(5.2), Inches(0.38), size=10, color=MI)
    ey += Inches(0.42)

# Right — frequency chart
card(sl, Inches(6.8), Inches(1.75), Inches(6.0), Inches(5.4))
rect(sl, Inches(6.8), Inches(1.75), Inches(6.0), Inches(0.04), fill=GO)
txb(sl, "BaseBJJ · Frequência — João Silva", Inches(7.0), Inches(1.88), Inches(5.5), Inches(0.28),
    size=8, color=GO, font="Courier New")

months_data = [("Jan",0), ("Fev",8), ("Mar",12), ("Abr",10), ("Mai",4)]
bar_max = 12
bar_area_h = Inches(1.8)
bar_w = Inches(0.75)
bar_start_x = Inches(7.1)
bar_base_y = Inches(4.1)

for i, (month, count) in enumerate(months_data):
    bx = bar_start_x + i * Inches(1.0)
    bar_h = (count / bar_max) * bar_area_h if count > 0 else Inches(0.04)
    bar_color = GO if i == 4 else RGBColor(0x30, 0x2C, 0x14)
    rect(sl, bx, bar_base_y - bar_h, bar_w, bar_h, fill=bar_color)
    txb(sl, str(count), bx, bar_base_y - bar_h - Inches(0.28), bar_w, Inches(0.25),
        size=12, bold=True, color=bar_color, font="Arial Black", align=PP_ALIGN.CENTER)
    txb(sl, month, bx, bar_base_y + Inches(0.04), bar_w, Inches(0.25),
        size=8, color=MI, font="Courier New", align=PP_ALIGN.CENTER)

rect(sl, Inches(6.9), bar_base_y, Inches(5.8), Inches(0.02), fill=DM)

freq_rows = [
    ("AULAS ESTE MÊS", "4", GO),
    ("TOTAL ACUMULADO", "34 aulas", WH),
    ("META MENSAL", "12 aulas", WH),
    ("PRÓX. AVALIAÇÃO", "Pendente →", GO),
]
fr_y = Inches(4.3)
for lbl, val, col in freq_rows:
    txb(sl, lbl, Inches(6.95), fr_y, Inches(3.2), Inches(0.28), size=8, color=MI, font="Courier New")
    txb(sl, val, Inches(10.2), fr_y, Inches(2.5), Inches(0.28), size=9, color=col, font="Courier New")
    rect(sl, Inches(6.9), fr_y + Inches(0.3), Inches(5.8), Inches(0.01), fill=DM)
    fr_y += Inches(0.34)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — GRADUAÇÕES
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 10)
txb(sl, "08", Inches(9.5), Inches(0.3), Inches(3.5), Inches(2.5),
    size=150, bold=True, color=RGBColor(0x15,0x12,0x05), font="Arial Black")

tag_label(sl, "ETAPA 08 · EVOLUÇÃO CONTÍNUA", Inches(0.55), Inches(0.55))
big_title(sl, "GESTÃO DE GRADUAÇÕES", Inches(0.55), Inches(0.85), size=42)
gold_bar(sl, Inches(0.55), Inches(1.65))

txb(sl, "CRITÉRIOS DE GRADUAÇÃO", Inches(0.55), Inches(1.82), Inches(5.5), Inches(0.28),
    size=8, color=MI, font="Courier New")
criteria = [
    ("A", "Frequência Mínima",
     "Aulas registradas no BaseBJJ por faixa. Rastreado automaticamente."),
    ("B", "Evolução Técnica",
     "Avaliação do professor: qualidade nas posições fundamentais da faixa atual."),
    ("C", "Comportamento e Atitude",
     "Respeito ao professor, companheiros e às regras do tatame. Disciplina."),
    ("D", "Tempo Mínimo na Faixa",
     "Diretrizes IBJJF respeitadas. Critérios transparentes para aluno e professor."),
]
sy = Inches(2.1)
for letter, title, desc in criteria:
    txb(sl, letter, Inches(0.55), sy, Inches(0.3), Inches(0.35),
        size=20, bold=True, color=RGBColor(0x44,0x3A,0x18), font="Arial Black")
    rect(sl, Inches(0.95), sy + Inches(0.06), Inches(5.0), Inches(0.015), fill=DM)
    txb(sl, title, Inches(0.95), sy + Inches(0.08), Inches(5.0), Inches(0.28),
        size=11, bold=True, color=WH)
    txb(sl, desc, Inches(0.95), sy + Inches(0.34), Inches(5.0), Inches(0.3),
        size=10, color=MI)
    sy += Inches(0.8)

card(sl, Inches(0.55), sy + Inches(0.1), Inches(5.5), Inches(1.1),
     gold_left=True, gold_bg=True)
txb(sl, "CERIMÔNIA DE GRADUAÇÃO", Inches(0.75), sy + Inches(0.22), Inches(5.2), Inches(0.25),
    size=7, color=GO, font="Courier New")
ceremony = [
    "Realizada no tatame, na frente da turma — reconhecimento público.",
    "Registro fotográfico e publicação no Instagram da academia.",
    "Faixa e data atualizadas no BaseBJJ pelo professor imediatamente.",
]
cy = sy + Inches(0.48)
for c in ceremony:
    rect(sl, Inches(0.75), cy + Inches(0.05), Inches(0.06), Inches(0.06), fill=GO)
    txb(sl, c, Inches(0.92), cy, Inches(5.2), Inches(0.28), size=10, color=WH)
    cy += Inches(0.28)

# Right — belt progression
txb(sl, "PROGRESSÃO DE FAIXAS — ADULTO", Inches(6.8), Inches(1.82), Inches(6.0), Inches(0.28),
    size=8, color=MI, font="Courier New")

by2 = Inches(2.1)
for bcolor, bname, bdesc, bb in BELTS:
    card(sl, Inches(6.8), by2, Inches(5.9), Inches(0.48))
    belt_bar_h(sl, Inches(6.95), by2 + Inches(0.13), Inches(0.48), Inches(0.22),
               bcolor, border_gold=bb)
    txb(sl, bname, Inches(7.58), by2 + Inches(0.05), Inches(3.5), Inches(0.25),
        size=10, bold=True, color=WH, font="Courier New")
    txb(sl, bdesc, Inches(7.58), by2 + Inches(0.26), Inches(3.5), Inches(0.25),
        size=9, color=MI)
    by2 += Inches(0.54)

# BaseBJJ graduation record
rect(sl, Inches(6.8), by2 + Inches(0.15), Inches(5.9), Inches(0.025), fill=DM)
card(sl, Inches(6.8), by2 + Inches(0.25), Inches(5.9), Inches(1.9))
rect(sl, Inches(6.8), by2 + Inches(0.25), Inches(5.9), Inches(0.04), fill=GO)
txb(sl, "BaseBJJ · Registro de Graduação", Inches(7.0), by2 + Inches(0.38), Inches(5.5), Inches(0.28),
    size=8, color=GO, font="Courier New")
grad_rows = [
    ("ALUNO",         "João Silva",    WH),
    ("FAIXA ANTERIOR","Branca · 3 graus", WH),
    ("NOVA FAIXA",    "Azul · 0 graus", RGBColor(0x25,0x63,0xEB)),
    ("DATA",          "15/05/2025",    WH),
    ("PROFESSOR",     "✓ Registrado",  GO),
]
gr_y = by2 + Inches(0.7)
for lbl, val, col in grad_rows:
    txb(sl, lbl, Inches(6.95), gr_y, Inches(2.5), Inches(0.26), size=8, color=MI, font="Courier New")
    txb(sl, val, Inches(9.5), gr_y, Inches(3.1), Inches(0.26), size=9, color=col, font="Courier New")
    rect(sl, Inches(6.9), gr_y + Inches(0.28), Inches(5.8), Inches(0.01), fill=DM)
    gr_y += Inches(0.3)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ENCERRAMENTO
# ═══════════════════════════════════════════════════════════════════════════════
sl = add_slide()
progress_bar(sl, 11)

tag_label(sl, "O CICLO NUNCA TERMINA", Inches(0.55), Inches(0.55))
big_title(sl, "O ALUNO NUNCA CHEGA AO\nFIM DA JORNADA", Inches(0.55), Inches(0.82), size=38)
gold_bar(sl, Inches(0.55), Inches(1.9))

# Final flow chart (compact)
flow_y2 = Inches(2.08)
cx_start = Inches(0.35)
step_w2 = Inches(1.35)
gap2 = Inches(0.06)
for i, (sid, short, lbl) in enumerate(FLOW_STEPS):
    sx = cx_start + i * (step_w2 + gap2)
    step_circle(sl, sx + Inches(0.42), flow_y2, short, done=True)
    txb(sl, lbl, sx, flow_y2 + Inches(0.54), step_w2, Inches(0.44),
        size=7.5, color=RGBColor(0x66,0x60,0x3A), font="Courier New", align=PP_ALIGN.CENTER)
    if i < len(FLOW_STEPS) - 1:
        arrow_right(sl, sx + step_w2 - Inches(0.02), flow_y2 + Inches(0.17))

# Loop arrow
lx = cx_start + len(FLOW_STEPS) * (step_w2 + gap2)
arrow_right(sl, lx - Inches(0.2), flow_y2 + Inches(0.17), w=Inches(0.22))
step_circle(sl, lx + Inches(0.25), flow_y2, "↺", done=False, hot=True)
txb(sl, "Reinicia\nCiclo", lx + Inches(0.08), flow_y2 + Inches(0.54), Inches(1.0), Inches(0.44),
    size=7.5, color=GO, font="Courier New", align=PP_ALIGN.CENTER)

rect(sl, Inches(0.35), Inches(3.2), Inches(12.7), Inches(0.015), fill=DM)

# 4 closing cards
closing = [
    ("RETENÇÃO NATURAL",   True,
     "A faixa seguinte cria um novo objetivo.\nQuem recebe um grau quer o próximo.\nA jornada nunca termina."),
    ("INDICAÇÃO",          False,
     "Alunos satisfeitos trazem amigos e\nfamiliares. O ciclo recomeça com\ncusto zero de aquisição."),
    ("DADOS & GESTÃO",     False,
     "Todo histórico no BaseBJJ: frequência,\npagamentos e graduações — decisões\nbaseadas em dados reais."),
    ("ECOSSISTEMA COMPLETO", True,
     "Site + BaseBJJ + WhatsApp formam um\nsistema integrado de captação,\nconversão e retenção."),
]
cw = Inches(3.0)
for i, (title, gold_bg, desc) in enumerate(closing):
    cx = Inches(0.35) + i * (cw + Inches(0.15))
    card(sl, cx, Inches(3.35), cw, Inches(2.5), gold_left=gold_bg, gold_bg=gold_bg)
    txb(sl, title, cx + Inches(0.2), Inches(3.48), cw - Inches(0.1), Inches(0.28),
        size=8, color=GO if gold_bg else MI, font="Courier New")
    txb(sl, desc, cx + Inches(0.2), Inches(3.82), cw - Inches(0.3), Inches(1.4),
        size=11, color=WH)

# Belt color bar footer
belt_colors_f = [
    RGBColor(0xF5,0xF5,0xF0), RGBColor(0x25,0x63,0xEB),
    RGBColor(0x7C,0x3A,0xED), RGBColor(0x92,0x40,0x0E), GO,
]
bw_f = Inches(2.6)
for i, bc in enumerate(belt_colors_f):
    rect(sl, Inches(0.35) + i * (bw_f + Inches(0.02)), Inches(6.35), bw_f, Inches(0.1),
         fill=bc if i > 0 else RGBColor(0x55,0x55,0x50))

txb(sl, "BLACKBOX.", Inches(0.55), Inches(6.58), Inches(4), Inches(0.4),
    size=20, bold=True, color=RGBColor(0x30,0x2C,0x14), font="Arial Black")
txb(sl, "APRESENTAÇÃO INTERNA · 2025",
    Inches(8.5), Inches(6.62), Inches(4.5), Inches(0.3),
    size=8, color=RGBColor(0x3A,0x38,0x30), font="Courier New", align=PP_ALIGN.RIGHT)

# ── SAVE ──────────────────────────────────────────────────────────────────────
output = "/home/user/arch-blueprint-revive/jornada-aluno-blackbox.pptx"
prs.save(output)
print(f"✓ Salvo em: {output}")
print(f"  Slides: {len(prs.slides)}")
